import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors (Excellence Team)
GREEN = RGBColor(0x1A, 0x43, 0x3A)       # #1A433A Deep Forest Green
ORANGE = RGBColor(0xFF, 0x70, 0x1A)      # #FF701A Tech Orange Accent
LIGHT_BG = RGBColor(0xF2, 0xF1, 0xED)    # #F2F1ED Soft Light Concrete
LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)  # Pure White
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)     # #0F172A Deep Tech Slate
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B Card Slate
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B
TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC
TEXT_MUTED_LIGHT = RGBColor(0x64, 0x74, 0x8B)
TEXT_MUTED_DARK = RGBColor(0x94, 0xA3, 0xB8)
BORDER_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

def add_header(slide, title_text, category_text, is_dark=False):
    # Category / Tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ORANGE
    p_cat.font.name = 'Consolas'

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT if is_dark else GREEN
    p_title.font.name = 'Arial'

def add_image_placeholder(slide, left, top, width, height, label, description, is_dark=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    if is_dark:
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        shape.line.color.rgb = ORANGE
    else:
        shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = f"🖼️ [ PLACEHOLDER IMAGE ]\n{label}"
    p1.font.bold = True
    p1.font.size = Pt(13)
    p1.font.color.rgb = ORANGE if is_dark else GREEN
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Consolas'

    p2 = tf.add_paragraph()
    p2.text = f"\nDescription visuelle :\n{description}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED_DARK if is_dark else TEXT_MUTED_LIGHT
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'

def add_footer(slide, is_dark=False):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "EXCELLENCE TEAM — PITCH DECK TEMPLATE  |  CONFIDENTIEL"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED_DARK if is_dark else TEXT_MUTED_LIGHT
    p.font.name = 'Consolas'

def build_presentation(is_dark=False):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg_color = DARK_BG if is_dark else LIGHT_BG
    card_bg = DARK_CARD if is_dark else LIGHT_CARD
    text_main = TEXT_LIGHT if is_dark else TEXT_DARK
    text_sub = TEXT_MUTED_DARK if is_dark else TEXT_MUTED_LIGHT

    # -------------------------------------------------------------
    # SLIDE 1 : HERO SLIDE / TITRE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = bg_color
    bg1.line.fill.background()

    # Accent bar left
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    # Title Box
    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(6.5), Inches(4.5))
    tf = tbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "EXCELLENCE TEAM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = 'Consolas'

    p2 = tf.add_paragraph()
    p2.text = "[ NOM DU PROJET ]"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = text_main
    p2.font.name = 'Arial'

    p3 = tf.add_paragraph()
    p3.text = "Sous-titre : Présentation stratégique et pitch deck opérationnel."
    p3.font.size = Pt(16)
    p3.font.color.rgb = text_sub
    p3.font.name = 'Arial'

    p4 = tf.add_paragraph()
    p4.text = "\n\nPrésenté par : [Nom / Équipe]\nDate : [DD/MM/YYYY]\nVersion : 1.0"
    p4.font.size = Pt(12)
    p4.font.color.rgb = text_sub
    p4.font.name = 'Consolas'

    add_image_placeholder(
        slide1, 8.0, 1.5, 4.5, 4.5,
        "HERO VISUAL / BRAND COVER",
        "Insérer l'image principale de marque, le logo 3D Brutaliste ou le mockup héro du projet.",
        is_dark
    )
    add_footer(slide1, is_dark)

    # -------------------------------------------------------------
    # SLIDE 2 : LE CONSTAT / PROBLÈME & OPPORTUNITÉ
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = bg_color
    bg2.line.fill.background()

    add_header(slide2, "Le Constat & L'Opportunité du Marché", "01 // LE DÉFI", is_dark)

    # 2 Text Cards
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = card_bg
    card1.line.color.rgb = GREEN if not is_dark else BORDER_GRAY
    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "⚠️ LE PROBLÈME IDENTIFIÉ"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ORANGE
    
    p = tf1.add_paragraph()
    p.text = "\n1. Absence de solutions agiles adaptées au marché local.\n2. Manque de passerelles concrètes entre jeunes talents et projets tech à fort impact.\n3. Friction d'exécution et lenteur de déploiement des outils existants."
    p.font.size = Pt(13)
    p.font.color.rgb = text_main

    add_image_placeholder(
        slide2, 6.7, 1.6, 5.8, 4.8,
        "VISUEL DU PROBLÈME / STATISTIQUES",
        "Graphique statistique, cartographie du marché ou illustration de la problématique.",
        is_dark
    )
    add_footer(slide2, is_dark)

    # -------------------------------------------------------------
    # SLIDE 3 : LA SOLUTION / PRODUIT (TECH & INNOVATION)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = bg_color
    bg3.line.fill.background()

    add_header(slide3, "Notre Solution : Innovation & Valeur Produit", "02 // LA SOLUTION", is_dark)

    # Product pitch text
    sol_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8))
    tf_s = sol_box.text_frame
    tf_s.word_wrap = True
    
    p = tf_s.paragraphs[0]
    p.text = "UNE APPROCHE RAW, TACTIQUE & PRECISENESS"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN if not is_dark else ORANGE

    p = tf_s.add_paragraph()
    p.text = "\nNotre plateforme apporte une réponse directe et scalable :"
    p.font.size = Pt(13)
    p.font.color.rgb = text_main

    bullets = [
        "🚀 Déploiement ultra-rapide et ergonomie pensée pour l'utilisateur.",
        "🔒 Sécurité renforcée et souveraineté des données.",
        "🤝 Conçu par une équipe jeune, dynamique et orientée résultat."
    ]
    for b in bullets:
        p = tf_s.add_paragraph()
        p.text = f"\n{b}"
        p.font.size = Pt(12)
        p.font.color.rgb = text_sub

    add_image_placeholder(
        slide3, 6.7, 1.6, 5.8, 4.8,
        "MOCKUP PRODUIT / DEMO UI",
        "Capture d'écran du Dashboard SaaS, de l'application mobile ou prototype fonctionnel.",
        is_dark
    )
    add_footer(slide3, is_dark)

    # -------------------------------------------------------------
    # SLIDE 4 : INFRASTRUCTURE & ARCHITECTURE (TECH FOCUS)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = bg_color
    bg4.line.fill.background()

    add_header(slide4, "Architecture Technique & Stack Technologique", "03 // TECH & BLUEPRINT", is_dark)

    # 3 Tech Cards horizontally
    card_w = Inches(3.6)
    gap = Inches(0.4)
    start_x = Inches(0.8)

    pillars = [
        ("CORE & DATA", "• Briques Open Source\n• Base de données scalable\n• API REST & GraphQL", GREEN),
        ("CLOUD & SÉCURITÉ", "• Hébergement hybride\n• Chiffrement de bout en bout\n• Monitoring temps réel", ORANGE),
        ("INTELLIGENCE / IA", "• Modèles IA optimisés\n• Automatisation des flux\n• Analyses prédictives", GREEN)
    ]

    for i, (p_title, p_desc, p_color) in enumerate(pillars):
        x = start_x + i * (card_w + gap)
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), card_w, Inches(2.2))
        c.fill.solid()
        c.fill.fore_color.rgb = card_bg
        c.line.color.rgb = p_color
        tf = c.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = p_title
        p0.font.bold = True
        p0.font.size = Pt(13)
        p0.font.color.rgb = p_color
        p0.font.name = 'Consolas'

        p1 = tf.add_paragraph()
        p1.text = f"\n{p_desc}"
        p1.font.size = Pt(11)
        p1.font.color.rgb = text_main

    add_image_placeholder(
        slide4, 0.8, 4.1, 11.7, 2.5,
        "BLUEPRINT ARCHITECTURE SYSTÈME",
        "Schéma d'architecture technique (Infrastructure Cloud, pipeline Data/IA et protocoles de sécurité).",
        is_dark
    )
    add_footer(slide4, is_dark)

    # -------------------------------------------------------------
    # SLIDE 5 : LA TEAM & VALEURS (TEAM FOCUS)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = bg_color
    bg5.line.fill.background()

    add_header(slide5, "La Team : Jeunesse, Vigueur & Excellence", "04 // HUMAIN & SYNERGIE", is_dark)

    # Team Pillars
    t_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = "NOS VALEURS FONDATRICES"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ORANGE
    p.font.name = 'Consolas'

    values = [
        ("HUMAN & SOLID", "Une équipe unie par le partage, l'entraide et l'amour de la technique."),
        ("TACTICAL & RAW", "Une capacité d'exécution brute, pragmatique et focalisée sur le terrain."),
        ("PREMIUM & AMBITIOUS", "Prouver que les étudiants créent une innovation de niveau mondial.")
    ]

    for title_v, desc_v in values:
        p = tf_t.add_paragraph()
        p.text = f"\n• {title_v}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = GREEN if not is_dark else ORANGE

        p = tf_t.add_paragraph()
        p.text = desc_v
        p.font.size = Pt(11)
        p.font.color.rgb = text_sub

    add_image_placeholder(
        slide5, 6.7, 1.6, 5.8, 4.8,
        "PHOTO TEAM / MOODBOARD COLLABORATIF",
        "Photo de l'équipe, moments de hackathon ou visuel illustrant l'énergie et la synergie de groupe.",
        is_dark
    )
    add_footer(slide5, is_dark)

    # -------------------------------------------------------------
    # SLIDE 6 : CONCLUSION & CTA (CLOSING)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = bg_color
    bg6.line.fill.background()

    # Big Closing Banner Box
    banner = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.7), Inches(4.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = card_bg
    banner.line.color.rgb = ORANGE
    banner.line.width = Pt(2)

    tf_b = banner.text_frame
    tf_b.word_wrap = True

    p = tf_b.paragraphs[0]
    p.text = "TEAM  •  TECH  •  INNOVATION"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ORANGE
    p.font.name = 'Consolas'

    p = tf_b.add_paragraph()
    p.text = "\nBâtissons l'Avenir Ensemble."
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_main

    p = tf_b.add_paragraph()
    p.text = "\nMontrer une équipe ancrée dans la tech, dont le but ultime est d'innover."
    p.font.size = Pt(14)
    p.font.color.rgb = text_sub

    p = tf_b.add_paragraph()
    p.text = "\n\n📍 Contact : contact@excellenceteam.com  |  🌐 Website : www.excellenceteam.com  |  💻 GitHub : @ExcellenceTeam"
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN if not is_dark else ORANGE
    p.font.name = 'Consolas'

    add_footer(slide6, is_dark)

    output_filename = "Excellence_Team_PitchDeck_DARK.pptx" if is_dark else "Excellence_Team_PitchDeck_LIGHT.pptx"
    prs.save(output_filename)
    print(f"Saved: {output_filename}")

if __name__ == "__main__":
    build_presentation(is_dark=False)
    build_presentation(is_dark=True)
